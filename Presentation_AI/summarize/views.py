
from django.shortcuts import render
from .forms import KeyWordForm
from django.http import HttpResponse
import wikipedia
import wikipediaapi
from congif.settings import WIKI_API_KEY, GOOGLE_API_KEY
import google.generativeai as genai
import json
from pptx import Presentation
from io import BytesIO
from PIL import Image
import os

def fetch_wikipedia_data(Key , top_n=1):
    wiki_wiki = wikipediaapi.Wikipedia(WIKI_API_KEY,'en')
    search_results = wikipedia.search(Key)
    
    articles = []
    for title in search_results[:top_n]:
        page = wiki_wiki.page(title)
        if page.exists():
            articles.append(page.text)
    return articles[:100]

def fetch_summarize_data(data,number):
    genai.configure(api_key=GOOGLE_API_KEY)
    model = genai.GenerativeModel('gemini-1.5-pro')
    response = model.generate_content(str(data)+"""by summarizing this data , Generate a JSON array with """+str(number)+""" objects in the following format:
        [
        {
            "title": "string",
            "content": "string",
            "image": "string"
        },
        ] where title is the title of silde , content is the content for the slide and image is suitable image description for the slide""")
    return response
def summarize(request):
    if request.method == 'POST':
        form = KeyWordForm(request.POST)
        if form.is_valid():
            key = form.cleaned_data['keyword']
            number = form.cleaned_data['number']
            text = fetch_wikipedia_data(Key=key)
            text = text[0][:5000]
            text = fetch_summarize_data(text,number).text
            String_json = ''
            holder = False
            for i in range(len(text)):
                if text[i] == '[': 
                    holder = True
                elif text[i] == ']':
                    String_json += ']'
                    holder = False
                if holder:
                    String_json += text[i]
            String_json = String_json.replace('\n', '')
            String_json = ' '.join(String_json.split())
            String_json = String_json.replace('`','').replace('} {','},{')
            data = json.loads(String_json)
            titles = [item['title'] for item in data]
            contents = [item['content'] for item in data]
            images = [item['image'] for item in data]
            
            return render(request , 'Data.html',{
                'form':form,
            })
    else:
        form = KeyWordForm()
    return render(request , 'Home.html',{
        'form':form
    })   
    
def create_presentation(titles, content):
    prs = Presentation()
    
    for title, text in zip(titles, content):
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        content_placeholder = slide.shapes.placeholders[1]
        title_placeholder.text = title
        content_placeholder.text = text
    
    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    
    return ppt_io